import io

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".html", ".htm", ".txt", ".md"}


def load_text(filename: str, content: bytes) -> str:
    """Extracts plain text from an uploaded document, dispatching on extension."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == ".pdf":
        return _load_pdf(content)
    if ext == ".docx":
        return _load_docx(content)
    if ext == ".pptx":
        return _load_pptx(content)
    if ext in (".xlsx", ".xls"):
        return _load_excel(content)
    if ext == ".csv":
        return _load_csv(content)
    if ext in (".html", ".htm"):
        return _load_html(content)
    # Plain text / markdown / anything else: decode as UTF-8 best-effort.
    return content.decode("utf-8", errors="ignore")


def _load_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _load_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def _load_pptx(content: bytes) -> str:
    """Sales decks/marketing materials (Phase 14a) are most often authored as
    .pptx — added specifically so those named asset types can be ingested in
    their native format, not just as a PDF export."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _load_csv(content: bytes) -> str:
    """pd.read_csv assumes one consistent-width table for the whole file —
    real exported documents (e.g. an invoice with a 2-column header block
    followed by a wider line-items table) are often ragged/multi-section and
    make it raise ParserError. Falls back to a plain UTF-8 decode (same
    fallback the "unknown extension" branch already uses) rather than ever
    failing the caller outright — a well-formed single-table CSV still gets
    the nicer to_string() tabulation unchanged."""
    try:
        return pd.read_csv(io.BytesIO(content)).to_string(index=False)
    except pd.errors.ParserError:
        return content.decode("utf-8", errors="ignore")


def _load_excel(content: bytes) -> str:
    sheets = pd.read_excel(io.BytesIO(content), sheet_name=None)
    return "\n\n".join(f"# {name}\n{df.to_string(index=False)}" for name, df in sheets.items())


def _load_html(content: bytes) -> str:
    soup = BeautifulSoup(content, "html.parser")
    return soup.get_text(separator="\n", strip=True)
